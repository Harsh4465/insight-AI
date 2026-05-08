import streamlit as st
import pandas as pd
from utils.viz_engine import render_hybrid_viz, generate_visual
from utils.ai_agent import get_executive_summary, generate_template_charts, get_ppt_storytelling
from utils.db_manager import load_charts_from_db, delete_chart_from_db, save_chart_to_db, supabase
import json
from pptx import Presentation
from pptx.util import Inches
import io
import plotly.io as pio

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_dark_theme(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # #0f172a

def apply_slide_branding(slide):
    apply_dark_theme(slide)
    left = Inches(0.5)
    top = Inches(7.1)
    width = Inches(9)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "INSIGHT AI | STRATEGIC COMMAND CENTER • CONFIDENTIAL"
    p.font.size = Inches(0.12)
    p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400

def generate_pptx_report(db_charts, df, chat_history_str=None):
    cols = df.columns.tolist() if df is not None else None
    story = get_ppt_storytelling(db_charts, chat_history_str, cols)
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_theme(slide)
    
    title = slide.shapes.title
    title.text = "Strategic Intelligence Report"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(99, 102, 241) 
    
    subtitle = slide.placeholders[1]
    subtitle.text = f"Cohesive Storytelling & Data Synthesis\nPowered by Insight AI Nexus"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(226, 232, 240)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_branding(slide)
    slide.shapes.title.text = "Executive Summary"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    body = slide.placeholders[1]
    body.text = str(story.get("exec_summary", "Strategic data mapping completed."))
    for p in body.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(226, 232, 240)

    for item in db_charts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        apply_slide_branding(slide)
        title_shape = slide.shapes.title
        title_shape.text = item['title']
        title_shape.text_frame.paragraphs[0].font.size = Inches(0.35)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        intent = item.get('intent', {})
        if isinstance(intent, str):
            try: intent = json.loads(intent)
            except: intent = {}

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "THE STORY:"
        p.font.bold = True
        p.font.color.rgb = RGBColor(34, 211, 238) 
        
        raw_story = story.get("slide_stories", {}).get(item['title'])
        if not raw_story:
            raw_story = item.get('ai_memo')
        if isinstance(raw_story, dict):
            raw_story = raw_story.get("answer", str(raw_story))
            
        # Ensure it's never empty on the slide
        if not raw_story or str(raw_story).strip() == "":
            raw_story = "Strategic Insight: This key performance indicator provides critical operational visibility. Further data accumulation will yield deeper predictive trends."
        
        p2 = tf.add_paragraph()
        p2.text = str(raw_story)
        p2.font.size = Inches(0.18)
        p2.font.color.rgb = RGBColor(226, 232, 240)

        if df is not None and intent:
            try:
                v_type, v_obj = generate_visual(intent, df)
                if v_type == "plotly":
                    img_bytes = pio.to_image(v_obj, format='png', width=800, height=600)
                    img_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(img_stream, Inches(4.8), Inches(1.8), width=Inches(4.7))
                elif v_type == "kpi":
                    val = v_obj.get("value", 0)
                    fval = f"{val:,.2f}" if isinstance(val, (int, float)) else str(val)
                    kpi_txBox = slide.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(4.0), Inches(2.0))
                    kpi_tf = kpi_txBox.text_frame
                    kp = kpi_tf.add_paragraph()
                    kp.text = fval
                    kp.font.size = Inches(0.8)
                    kp.font.bold = True
                    kp.font.color.rgb = RGBColor(56, 189, 248) # #38bdf8 light blue
                    kp.alignment = PP_ALIGN.CENTER
            except Exception as e:
                print(f"PPTX Image Error: {e}")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_branding(slide)
    slide.shapes.title.text = "Strategic Roadmap"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    body = slide.placeholders[1]
    recs = story.get("strategic_recommendations", ["Monitor trends", "Scale operations"])
    body.text = "\n".join([f"• {r}" for r in recs])
    for p in body.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(226, 232, 240)

    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    return ppt_output


def phase4_dashboard():
    st.markdown('<div class="fade-in">', unsafe_allow_html=True)
    
    col_h1, col_h2 = st.columns([3, 1])
    with col_h1:
        st.markdown("""
            <div style="margin-bottom: 1rem;">
                <h1 style="margin:0;">📉 Power <span class="text-gradient">Dashboard</span></h1>
                <p style="color: var(--text-dim); margin-top:0.5rem; font-size: 1.1rem;">Your customized strategic control center.</p>
            </div>
        """, unsafe_allow_html=True)
    
    user = st.session_state.get('user')
    if not user:
        st.warning("Please login to access your gallery.")
        return

    with st.expander("✨ Auto-Generate AI Templates", expanded=False):
        t1, t2, t3 = st.columns(3)
        if t1.button("📈 Executive Overview", use_container_width=True): st.session_state.trigger_template = "Executive Overview"
        if t2.button("🎯 Performance & Deep Dive", use_container_width=True): st.session_state.trigger_template = "Performance & Deep Dive"
        if t3.button("🌍 Demographics & Trends", use_container_width=True): st.session_state.trigger_template = "Demographics & Trends"

    if st.session_state.get("trigger_template"):
        temp = st.session_state.trigger_template
        st.session_state.trigger_template = None
        if st.session_state.df is not None:
            with st.spinner(f"Insight AI is building the {temp} Dashboard..."):
                intents = generate_template_charts(temp, st.session_state.df)
                if intents:
                    for intent in intents:
                        save_chart_to_db(user.id, intent.get('title', 'Generated Chart'), intent.get('type', 'bar'), intent)
                    st.toast("✅ Template applied successfully!")
                    st.rerun()
                else:
                    st.error("Failed to generate template.")
        else:
            st.warning("Please connect a dataset first.")

    db_charts = load_charts_from_db(user.id)
    if not db_charts:
        st.info("Your dashboard is empty. Select a Template above or pin insights from the Chat Hub.")
        return

    # --- 1. NATIVE CHART CROSS-FILTERING ---
    df = st.session_state.df
    filtered_df = df.copy() if df is not None else None
    
    if df is not None:
        active_filters = {}
        for item in db_charts:
            chart_key = f"dash_chart_{item['id']}"
            if chart_key in st.session_state:
                selection = st.session_state[chart_key]
                points = []
                
                # Handle both dict and object structures depending on Streamlit version
                if isinstance(selection, dict):
                    points = selection.get("selection", {}).get("points", [])
                elif hasattr(selection, "selection"):
                    raw_points = getattr(selection.selection, "points", [])
                    if isinstance(raw_points, list):
                        for p in raw_points:
                            if isinstance(p, dict): points.append(p)
                            elif hasattr(p, "x"): points.append({"x": p.x})

                if points:
                    x_vals = []
                    for p in points:
                        val = p.get("x")
                        if val is None: val = p.get("label") # Fallback for pie/funnel charts
                        if val is not None: x_vals.append(val)
                    
                    intent = item.get('intent', {})
                    if isinstance(intent, str):
                        try: intent = json.loads(intent)
                        except: intent = {}
                    x_col = intent.get("x")
                    
                    if x_col and x_vals and x_col in filtered_df.columns:
                        active_filters[x_col] = active_filters.get(x_col, []) + x_vals
        
        if active_filters:
            filter_desc = " | ".join([f"{k}: {', '.join(map(str, set(v)))}" for k, v in active_filters.items()])
            st.markdown(f"""
                <div style="background: rgba(16, 185, 129, 0.1); padding: 0.8rem; border-radius: 8px; border-left: 4px solid var(--success); margin-bottom: 2rem;">
                    <span style="color: var(--success); font-weight: bold;">⚡ Cross-Filter Active</span>
                    <span style="color: var(--text-dim); font-size: 0.9rem; margin-left: 10px;">Filtering [{filter_desc}]. Double-click chart to clear.</span>
                </div>
            """, unsafe_allow_html=True)
            for col, vals in active_filters.items():
                backup_df = filtered_df.copy()
                
                is_num = pd.api.types.is_numeric_dtype(filtered_df[col])
                if is_num:
                    try:
                        num_vals = [float(v) for v in vals]
                        filtered_df = filtered_df[filtered_df[col].apply(lambda x: any(abs(float(x) - v) < 1e-5 for v in num_vals) if pd.notnull(x) else False)]
                    except Exception:
                        str_vals = [str(v).strip() for v in vals]
                        filtered_df = filtered_df[filtered_df[col].astype(str).str.strip().isin(str_vals)]
                else:
                    str_vals = [str(v).strip() for v in vals]
                    filtered_df = filtered_df[filtered_df[col].astype(str).str.strip().isin(str_vals)]
                
                # Graceful degradation: If filter breaks everything (e.g. clicked a histogram bin), revert to prevent crash
                if filtered_df.empty:
                    filtered_df = backup_df

    # Split Intents
    kpi_items = []
    chart_items = []
    for item in db_charts:
        intent = item.get('intent', {})
        if isinstance(intent, str):
            try: intent = json.loads(intent)
            except: intent = {}
        item['parsed_intent'] = intent
        
        if intent.get("type", "").lower() == "kpi" or item.get('chart_type') == 'kpi':
            kpi_items.append(item)
        else:
            chart_items.append(item)

    with col_h2:
        if st.button("📊 Export Deck", use_container_width=True, type="secondary"):
            with st.spinner("Creating professional slides..."):
                chat_history_str = ""
                if "messages" in st.session_state:
                    chat_history_str = "\n".join([f"{m['role']}: {m['content'][:200]}" for m in st.session_state.messages[-10:]])
                ppt_data = generate_pptx_report(db_charts, filtered_df, chat_history_str)
                st.download_button(
                    label="📥 Download PPTX",
                    data=ppt_data,
                    file_name=f"Strategic_Report_{st.session_state.get('filename', 'export')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        
        if st.button("🌐 Export HTML", use_container_width=True, type="primary"):
            with st.spinner("Compiling interactive offline dashboard..."):
                safe_df = filtered_df.fillna("") if filtered_df is not None else pd.DataFrame()
                json_data_str = safe_df.to_json(orient="records")
                
                html_content = f"""
                <html>
                <head>
                    <title>Insight AI Dashboard</title>
                    <style>
                        body {{ background: #0f172a; color: white; font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif; padding: 20px; margin: 0; }}
                        .header {{ text-align: left; margin-bottom: 20px; border-bottom: 1px solid #334155; padding-bottom: 10px; display: flex; justify-content: space-between; align-items: flex-end; }}
                        h1 {{ color: #e2e8f0; margin: 0; font-size: 28px; font-weight: 600; }}
                        .subtitle {{ color: #94a3b8; font-size: 14px; margin-top: 5px; }}
                        
                        .kpi-container {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 15px; margin-bottom: 25px; }}
                        .kpi-card {{ background: #1e293b; border-left: 4px solid #38bdf8; border-radius: 8px; padding: 15px; box-shadow: 0 4px 6px rgba(0,0,0,0.2); transition: all 0.3s ease; }}
                        .kpi-title {{ font-size: 12px; color: #94a3b8; text-transform: uppercase; font-weight: 600; letter-spacing: 0.5px; }}
                        .kpi-val {{ font-size: 32px; font-weight: bold; color: #f8fafc; margin-top: 5px; }}
                        
                        .chart-grid {{ display: grid; grid-template-columns: repeat(2, 1fr); gap: 20px; }}
                        @media (max-width: 1000px) {{ .chart-grid {{ grid-template-columns: 1fr; }} }}
                        .chart-card {{ background: #1e293b; border: 1px solid #334155; border-radius: 12px; padding: 15px; box-shadow: 0 4px 10px rgba(0,0,0,0.3); }}
                        .chart-title {{ font-size: 16px; color: #f8fafc; margin-top: 0; margin-bottom: 15px; font-weight: 600; }}
                        .chart-wrapper {{ width: 100%; height: 350px; overflow: hidden; }}
                    </style>
                    <script src="https://cdn.plot.ly/plotly-2.32.0.min.js"></script>
                    <script>
                        window.dashData = {json_data_str};
                        window.originalDashData = window.dashData;
                        window.kpiConfig = [];

                        function calculateAggregate(data, col, agg) {{
                            if (!data || data.length === 0) return 0;
                            if (agg === 'count') return data.length;
                            
                            let sum = 0;
                            let count = 0;
                            data.forEach(row => {{
                                let val = parseFloat(row[col]);
                                if (!isNaN(val)) {{
                                    sum += val;
                                    count += 1;
                                }}
                            }});
                            if (agg === 'sum') return sum;
                            if (agg === 'mean') return count > 0 ? sum / count : 0;
                            return data.length;
                        }}

                        function updateKPIs(data) {{
                            window.kpiConfig.forEach(kpi => {{
                                let val = calculateAggregate(data, kpi.y, kpi.agg);
                                let el = document.getElementById(kpi.id);
                                if(el) {{
                                    el.innerText = val.toLocaleString(undefined, {{minimumFractionDigits: 2, maximumFractionDigits: 2}});
                                }}
                            }});
                        }}

                        function applyCrossFilter(xCol, xVal) {{
                            if(!xCol || !xVal) return;
                            let filtered = window.originalDashData.filter(row => row[xCol] == xVal);
                            updateKPIs(filtered);
                        }}
                        
                        function resetCrossFilter() {{
                            updateKPIs(window.originalDashData);
                        }}
                    </script>
                </head>
                <body>
                    <div class="header">
                        <div>
                            <h1>Insight AI Power Dashboard</h1>
                            <div class="subtitle">Data Source: {st.session_state.get('filename', 'Live Dataset')} | Interactive Offline Mode Enabled</div>
                        </div>
                        <div class="subtitle">Click any chart bar/slice to filter KPIs. Double-click background to reset.</div>
                    </div>
                    
                    <div class="kpi-container">
                """
                
                kpi_configs = []
                for i, item in enumerate(kpi_items):
                    kpi_id = f"kpi_{i}"
                    y_col = item['parsed_intent'].get('y', '')
                    agg = item['parsed_intent'].get('agg', 'sum')
                    kpi_configs.append({"id": kpi_id, "y": y_col, "agg": agg})
                    
                    html_content += f"""
                        <div class="kpi-card">
                            <div class="kpi-title">{item['title']}</div>
                    """
                    if filtered_df is not None:
                        v_type, v_obj = generate_visual(item['parsed_intent'], filtered_df)
                        if v_type == 'kpi' and isinstance(v_obj, dict):
                            val = v_obj.get("value", 0)
                            fval = f"{val:,.2f}" if isinstance(val, (int, float)) else str(val)
                            html_content += f'<div class="kpi-val" id="{kpi_id}">{fval}</div>'
                        else:
                            html_content += f'<div class="kpi-val" id="{kpi_id}">Err</div>'
                    html_content += "</div>"

                html_content += f"""
                    </div>
                    <script>
                        window.kpiConfig = {json.dumps(kpi_configs)};
                    </script>
                    <div class="chart-grid">
                """

                import uuid
                chart_scripts = ""
                for item in chart_items:
                    html_content += f"""
                        <div class="chart-card">
                            <h3 class="chart-title">{item['title']}</h3>
                            <div class="chart-wrapper">
                    """
                    if filtered_df is not None:
                        v_type, v_obj = generate_visual(item['parsed_intent'], filtered_df)
                        if v_type == "plotly" and v_obj:
                            chart_id = f"c_{uuid.uuid4().hex}"
                            x_col = item['parsed_intent'].get('x', '')
                            chart_html = v_obj.to_html(full_html=False, include_plotlyjs=False, div_id=chart_id)
                            html_content += chart_html
                            
                            chart_scripts += f"""
                                var plotEl_{chart_id} = document.getElementById('{chart_id}');
                                if(plotEl_{chart_id}) {{
                                    plotEl_{chart_id}.on('plotly_click', function(data){{
                                        if(data.points && data.points.length > 0) {{
                                            let xVal = data.points[0].x;
                                            applyCrossFilter('{x_col}', xVal);
                                        }}
                                    }});
                                    plotEl_{chart_id}.on('plotly_doubleclick', function(){{
                                        resetCrossFilter();
                                    }});
                                }}
                            """
                        else:
                            html_content += '<div style="color:#64748b; padding:20px;">Visualization not supported in static HTML.</div>'
                    html_content += """
                            </div>
                        </div>
                    """
                
                html_content += f"""
                    </div>
                    <script>
                        window.addEventListener('load', function() {{
                            {chart_scripts}
                        }});
                    </script>
                </body>
                </html>
                """

                st.download_button(
                    label="📥 Click to Download HTML",
                    data=html_content,
                    file_name=f"Power_Dashboard_{st.session_state.get('filename', 'export')}.html",
                    mime="text/html",
                    use_container_width=True
                )


    st.markdown("<br>", unsafe_allow_html=True)

    if kpi_items:
        cols_per_row = min(4, max(2, len(kpi_items)))
        kpi_cols = st.columns(cols_per_row)
        for i, item in enumerate(kpi_items):
            with kpi_cols[i % cols_per_row]:
                if item['parsed_intent'] and filtered_df is not None:
                    v_type, v_obj = generate_visual(item['parsed_intent'], filtered_df)
                    if v_obj:
                        render_hybrid_viz(v_type, v_obj, title=item['title'], show_pin=False, key=f"kpi_{i}")
                        if st.button("🗑️", key=f"rm_kpi_{item['id']}", help="Remove KPI", use_container_width=True):
                            if delete_chart_from_db(item['id']): st.rerun()
                    else:
                        st.warning("KPI Data Missing")
    
    st.markdown("<br><hr style='border-color: rgba(255,255,255,0.1);'><br>", unsafe_allow_html=True)

    if chart_items:
        for i in range(0, len(chart_items), 2):
            cols = st.columns(2)
            for j in range(2):
                if i + j < len(chart_items):
                    item = chart_items[i + j]
                    intent = item['parsed_intent']
                    
                    with cols[j]:
                        st.markdown(f'<div class="glass-card" style="margin-bottom: 1.5rem; padding: 1.5rem;">', unsafe_allow_html=True)
                        
                        tc1, tc2 = st.columns([5, 1])
                        with tc1:
                            st.markdown(f"<h3 style='margin:0; font-size:1.1rem; white-space:nowrap; overflow:hidden; text-overflow:ellipsis;'>{item['title']}</h3>", unsafe_allow_html=True)
                        with tc2:
                            if st.button("✖", key=f"rm_chart_{item['id']}", help="Remove Chart"):
                                if delete_chart_from_db(item['id']): st.rerun()

                        if intent and df is not None:
                            v_type, v_obj = generate_visual(intent, df)
                            if v_obj:
                                render_hybrid_viz(v_type, v_obj, show_pin=False, key=f"dash_chart_{item['id']}")
                            else:
                                st.warning("Chart render failed.")
                        else:
                            st.info("⚠️ Re-upload dataset to view chart.")

                        current_memo = item.get('ai_memo', '') or ""
                        with st.expander("🪄 AI Narrative & Analysis"):
                            if current_memo:
                                st.markdown(f"<div style='font-size:0.9rem; color:var(--text-dim); line-height:1.5;'>{current_memo}</div>", unsafe_allow_html=True)
                                try:
                                    x, y = intent.get('x'), intent.get('y')
                                    if x and y and x in filtered_df.columns and y in filtered_df.columns:
                                        csv_data = filtered_df[[x, y]].head(100).to_csv(index=False)
                                        st.download_button("📥 Download CSV", data=csv_data, file_name=f"{item['title']}.csv", mime="text/csv", key=f"dl_{item['id']}", use_container_width=True)
                                except: pass
                            else:
                                if st.button("✨ Generate AI Analysis", key=f"ai_gen_{item['id']}", use_container_width=True):
                                    with st.spinner("Analyzing..."):
                                        summary = get_executive_summary(item['title'], intent, filtered_df)
                                        if summary:
                                            supabase.table("saved_charts").update({"ai_memo": summary}).eq("id", item['id']).execute()
                                            st.rerun()

                        st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("<br><br>", unsafe_allow_html=True)
    if st.button("🏠 Return to Home Hub", type="secondary", use_container_width=True):
        st.session_state.current_page = "Home"
        st.rerun()

    st.markdown('</div>', unsafe_allow_html=True)
